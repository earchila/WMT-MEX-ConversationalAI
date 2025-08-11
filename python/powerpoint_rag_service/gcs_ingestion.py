import os
from google.cloud import storage
from google.cloud import aiplatform
from pptx_reader import extract_text_from_pptx
import json

def download_pptx_from_gcs(bucket_name: str, blob_name: str, destination_file_name: str):
    """Downloads a blob from the bucket."""
    storage_client = storage.Client()
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(blob_name)
    blob.download_to_filename(destination_file_name)
    print(f"Blob {blob_name} downloaded to {destination_file_name}.")

def upload_to_gcs(bucket_name: str, source_file_name: str, destination_blob_name: str):
    """Uploads a file to the bucket."""
    storage_client = storage.Client()
    bucket = storage_client.bucket(bucket_name)
    blob = bucket.blob(destination_blob_name)
    blob.upload_from_filename(source_file_name)
    print(f"File {source_file_name} uploaded to {destination_blob_name}.")

def generate_embeddings(texts: list[str], project_id: str, location: str) -> list[list[float]]:
    """Generates embeddings for a list of texts using Vertex AI Embedding Model."""
    aiplatform.init(project=project_id, location=location)
    model = aiplatform.TextEmbeddingModel.from_pretrained("text-embedding-004")
    embeddings = model.predict(texts)
    return [embedding.values for embedding in embeddings]

def process_pptx_for_vector_search(
    bucket_name: str,
    pptx_blob_name: str,
    project_id: str,
    location: str,
    output_gcs_path: str, # e.g., "gs://your-bucket/vector_search_data/"
):
    """
    Processes a PowerPoint file from GCS, extracts text, generates embeddings,
    and uploads the data to GCS in a format suitable for Vertex AI Vector Search.
    """
    local_pptx_path = f"/tmp/{os.path.basename(pptx_blob_name)}"
    download_pptx_from_gcs(bucket_name, pptx_blob_name, local_pptx_path)

    extracted_text = extract_text_from_pptx(local_pptx_path)
    
    # Split text into chunks if it's too long for embedding model or for better retrieval
    # For simplicity, we'll treat each slide's text as a chunk for now.
    # A more sophisticated approach would involve sentence splitting and chunking.
    
    # Assuming extract_text_from_pptx returns text per slide or a single string
    # If it returns a single string, we might want to split it into paragraphs/sentences
    
    # For now, let's assume extracted_text is a single string, and we'll embed it as one.
    # In a real scenario, you'd chunk this text.
    
    if not extracted_text:
        print(f"No text extracted from {pptx_blob_name}. Skipping embedding.")
        return

    # Generate embeddings for the extracted text
    # Note: text-embedding-004 has a limit of 3072 input tokens per request.
    # For very large PPTX files, you'd need to chunk the text and embed chunks.
    embeddings = generate_embeddings([extracted_text], project_id, location)[0] # Get the first (and only) embedding

    # Prepare data in JSONL format for Vertex AI Vector Search
    # Each line is a JSON object with 'id', 'embedding', and optionally 'restricts' and 'numeric_restricts'
    # 'id' should be unique for each data point.
    
    # Using the blob name as a unique ID for simplicity
    data_point = {
        "id": pptx_blob_name.replace("/", "_").replace(".", "_"), # Create a valid ID
        "embedding": embeddings,
        "metadata": {"source_file": pptx_blob_name, "text_content": extracted_text}
    }

    output_file_name = f"{os.path.basename(pptx_blob_name)}.jsonl"
    local_output_path = f"/tmp/{output_file_name}"
    
    with open(local_output_path, "w") as f:
        f.write(json.dumps(data_point) + "\n")

    # Upload the JSONL file to the specified GCS path
    destination_blob_name = f"{output_gcs_path.replace('gs://', '')}{output_file_name}"
    upload_to_gcs(bucket_name, local_output_path, destination_blob_name)

    os.remove(local_pptx_path)
    os.remove(local_output_path)
    print(f"Successfully processed and uploaded data for {pptx_blob_name}")

if __name__ == "__main__":
    # Example Usage (replace with your actual values)
    # Ensure you have authenticated to GCP (e.g., `gcloud auth application-default login`)
    # and enabled the necessary APIs (Vertex AI API, Cloud Storage API).

    # Dummy values for demonstration
    GCS_BUCKET = os.getenv("GCS_POWERPOINT_BUCKET", "your-powerpoint-bucket")
    GCP_PROJECT = os.getenv("GOOGLE_CLOUD_PROJECT", "your-gcp-project-id")
    GCP_LOCATION = os.getenv("GOOGLE_CLOUD_LOCATION", "us-central1")
    
    # Create a dummy PowerPoint file for testing ingestion
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    title.text = "Sample Presentation for RAG"
    body.text = "This is a test slide for the RAG system.\nIt contains some important information."
    
    dummy_pptx_name = "test_rag_presentation.pptx"
    prs.save(dummy_pptx_name)
    
    # Upload dummy PPTX to GCS
    upload_to_gcs(GCS_BUCKET, dummy_pptx_name, dummy_pptx_name)
    
    # Process the uploaded PPTX
    process_pptx_for_vector_search(
        bucket_name=GCS_BUCKET,
        pptx_blob_name=dummy_pptx_name,
        project_id=GCP_PROJECT,
        location=GCP_LOCATION,
        output_gcs_path=f"gs://{GCS_BUCKET}/vector_search_data/"
    )
    
    os.remove(dummy_pptx_name)