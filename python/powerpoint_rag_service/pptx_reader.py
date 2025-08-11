from pptx import Presentation

def extract_text_from_pptx(pptx_path):
    """
    Extracts all text from a PowerPoint presentation.
    """
    text_content = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
    except Exception as e:
        return f"Error reading PowerPoint file: {e}"
    return "\n".join(text_content)

if __name__ == "__main__":
    # Example usage:
    # Create a dummy PowerPoint file for testing
    from pptx.util import Inches
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Test Presentation"
    body.text = "This is the first slide.\nIt contains some text."
    
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.placeholders[1]
    
    title.text = "Second Slide"
    body.text = "More content here.\nIncluding a second paragraph."
    
    dummy_pptx_path = "dummy_presentation.pptx"
    prs.save(dummy_pptx_path)
    
    print(f"Created dummy PowerPoint file: {dummy_pptx_path}")
    
    extracted_text = extract_text_from_pptx(dummy_pptx_path)
    print("\nExtracted Text:")
    print(extracted_text)